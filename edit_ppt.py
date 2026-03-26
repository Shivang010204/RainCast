from pptx import Presentation

def update_ppt(input_path, output_path):
    prs = Presentation(input_path)

    replacements = {
        # Slide 6
        "Hybrid Intelligence:": "Smart Predictions: Brings together AI and real user reports to predict weather better.\nCustom Advice: Gives specific safety tips for Farmers and Construction Workers based on the weather.\nPhoto Check: When users report rain, they must upload a photo. The system checks the hidden details (metadata) to make sure it's real and fresh.\nCommunity Votes: Users can vote on weather reports to build trust in the community.\nRemembers You: Saves your city searches and preferences as you use the app.",
        
        # Slide 7
        "Weather Retrieval Module:": "Weather Fetcher: Gets live weather data like Temperature, Humidity, and Wind from the internet.\nAI Brain: Uses smart formulas to look at local data and guess if it will rain.\nPhoto Verifier: Asks users to upload a photo to report rain, and checks the photo details to stop fake reports.\nAdmin Panel: A safe page for admins to watch over reports, check user trust, and delete bad posts or images.\nAdvice Giver: Creates helpful 'Dos and Don'ts' depending on whether you are a Farmer or Construction Worker.",
        
        # Slide 8
        "Metadata Validation:": "Photo Checks: Rejects fake or old photos automatically.\nDisagreement Alerts: Warns you if the AI says 'No Rain' but a user just reported 'Rain'.\nModern Look: A clean and cool dark design that is easy to read outside in the sun.\nSaving History: Keeps a record of every search and report to help us study the weather over time.",

        # Slide 9
        "Current Scope:": "What we do Now: Helps Farmers and Construction Workers around the world make choices based on the weather.\nWhat we will do Later:\n- Send SMS text messages to warn farmers about sudden rain.\n- Use GPS to find your location automatically.\n- Add more languages to help people in rural areas.",

        # Slide 13
        "Process Model: Iterative Model": "How We Built It: Step by Step. We made the smart AI brain first, then added user reporting and admin tools by testing and improving.\nHow It's Organized: It's built in separate parts that talk to each other, like getting weather data, running the AI, and checking photos.\nWhat Kind of App It Is: It's an app that helps you make decisions and uses reports from the community.",

        # Slide 16
        "Database: We use a CSV": "Database: We keep our records simple, like a digital notepad. Every search adds a new line with the time, city, and weather.\nScreens: The app looks modern, like frosted glass. It has your controls on the left, weather in the middle, and maps on the right.\nOver the Day: A chart shows you what the temperature will be for the next 8 hours so you know what's coming.",

        # Slide 17 (Testing)
        "Search Test: Typed": "Basic Tests:\n- City Search: Typed a city -> Map moved, and weather showed up. (Success)\n- Role Switch: Picked 'Farmer' or 'Builder' -> Got the right advice. (Success)\n\nSecurity Tests:\n- Fake Photo Test: Uploaded an old photo -> App blocked it. (Success)\n- Spam Report Test: Tried to send too many reports -> App asked me to wait. (Success)\n\nReal World Tests:\n- Mobile Phone Check: Used a phone browser -> Fast and easy to read. (Success)\n- Location Match: Checked the weather for a rainy city -> AI agreed it was raining. (Success)\n- Dark Mode Test: Looked at the screen outside in the sun -> Very clear to read. (Success)"
    }

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                for key, replacement in replacements.items():
                    if key in shape.text:
                        shape.text = replacement
                        # Also adjust font size slightly down since we added text in some places
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.font:
                                    # Optional: Can try to set smaller font size if needed
                                    pass

    prs.save(output_path)
    print(f"Saved simplified presentation to {output_path}")

if __name__ == '__main__':
    update_ppt("raincasts.pptx", "raincasts_simplified.pptx")
