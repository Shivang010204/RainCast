import sys
import collections
import collections.abc
from pptx import Presentation

def extract_text(ppt_path):
    try:
        prs = Presentation(ppt_path)
    except Exception as e:
        print(f"Error opening {ppt_path}: {e}")
        return

    for i, slide in enumerate(prs.slides):
        print(f"--- Slide {i+1} ---")
        for j, shape in enumerate(slide.shapes):
            if hasattr(shape, "text") and shape.text.strip():
                print(f"Shape {j}: {repr(shape.text)}")
        print("\n")

if __name__ == '__main__':
    extract_text("raincasts.pptx")
